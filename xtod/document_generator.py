#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档生成器 - 支持 PDF 和 PPT 格式
"""

import os
from datetime import datetime
from typing import Dict, List


class PDFGenerator:
    """PDF 文档生成器"""

    def __init__(self):
        """初始化 PDF 生成器"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.units import inch
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
            from reportlab.lib import colors
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            self.A4 = A4
            self.inch = inch
            self.getSampleStyleSheet = getSampleStyleSheet
            self.ParagraphStyle = ParagraphStyle
            self.TA_LEFT = TA_LEFT
            self.TA_CENTER = TA_CENTER
            self.SimpleDocTemplate = SimpleDocTemplate
            self.Paragraph = Paragraph
            self.Spacer = Spacer
            self.Image = Image
            self.PageBreak = PageBreak
            self.colors = colors
            self.pdfmetrics = pdfmetrics
            self.TTFont = TTFont

            self._load_chinese_font()
        except ImportError:
            print("⚠️ reportlab 未安装，正在安装...")
            import subprocess
            subprocess.run(['pip3', 'install', 'reportlab'], check=True)
            # 重新导入
            self.__init__()

    def _load_chinese_font(self):
        """加载中文字体"""
        try:
            font_paths = [
                '/System/Library/Fonts/STHeiti Light.ttc',
                '/System/Library/Fonts/PingFang.ttc',
                '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'
            ]

            for font_path in font_paths:
                if os.path.exists(font_path):
                    self.pdfmetrics.registerFont(self.TTFont('ChineseFont', font_path))
                    self.chinese_font = 'ChineseFont'
                    print(f"✅ 已加载字体: {font_path}")
                    return

            self.chinese_font = 'Helvetica'
            print("⚠️ 未找到中文字体，使用 Helvetica")
        except Exception as e:
            self.chinese_font = 'Helvetica'
            print(f"⚠️ 字体加载失败: {e}")

    def generate(self, data: Dict, output_path: str) -> str:
        """
        生成 PDF 文档

        Args:
            data: Twitter thread 数据
            output_path: 输出文件路径

        Returns:
            生成的文件路径
        """
        print("📄 正在生成 PDF...")

        # 创建 PDF
        doc = self.SimpleDocTemplate(
            output_path,
            pagesize=self.A4,
            rightMargin=0.75 * self.inch,
            leftMargin=0.75 * self.inch,
            topMargin=1 * self.inch,
            bottomMargin=0.75 * self.inch
        )

        # 定义样式
        styles = self.getSampleStyleSheet()

        title_style = self.ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName=self.chinese_font,
            fontSize=24,
            textColor=self.colors.HexColor('#1DA1F2'),
            spaceAfter=20,
            alignment=self.TA_CENTER
        )

        subtitle_style = self.ParagraphStyle(
            'Subtitle',
            parent=styles['Normal'],
            fontName=self.chinese_font,
            fontSize=14,
            textColor=self.colors.HexColor('#14171A'),
            spaceAfter=10,
            alignment=self.TA_CENTER
        )

        content_style = self.ParagraphStyle(
            'Content',
            parent=styles['Normal'],
            fontName=self.chinese_font,
            fontSize=11,
            textColor=self.colors.HexColor('#14171A'),
            spaceAfter=12,
            leading=16,
            alignment=self.TA_LEFT
        )

        time_style = self.ParagraphStyle(
            'Time',
            parent=styles['Normal'],
            fontName=self.chinese_font,
            fontSize=9,
            textColor=self.colors.HexColor('#657786'),
            spaceAfter=10
        )

        metrics_style = self.ParagraphStyle(
            'Metrics',
            parent=styles['Normal'],
            fontName=self.chinese_font,
            fontSize=9,
            textColor=self.colors.HexColor('#657786'),
            spaceAfter=15
        )

        # 构建内容
        story = []

        # 封面
        author = data.get('author', {})
        author_name = author.get('name', 'Twitter User')
        total = data.get('total_tweets', 0)

        story.append(self.Paragraph(f"{author_name} 的 Twitter Thread", title_style))
        story.append(self.Spacer(1, 0.3 * self.inch))
        story.append(self.Paragraph(f"共 {total} 条推文", subtitle_style))
        story.append(self.Spacer(1, 0.5 * self.inch))
        story.append(self.Paragraph(
            f"生成时间：{datetime.now().strftime('%Y年%m月%d日 %H:%M')}",
            time_style
        ))
        story.append(self.PageBreak())

        # 添加每条推文
        for i, tweet in enumerate(data.get('tweets', []), 1):
            # 推文标题
            time_text = tweet.get('time_text', '')
            header = f"<b>推文 #{i}</b> · {time_text}"
            story.append(self.Paragraph(header, subtitle_style))
            story.append(self.Spacer(1, 0.1 * self.inch))

            # 推文截图
            screenshot = tweet.get('screenshot')
            if screenshot and os.path.exists(screenshot):
                try:
                    img = self.Image(screenshot, width=5*self.inch, height=3.5*self.inch)
                    story.append(img)
                    story.append(self.Spacer(1, 0.2 * self.inch))
                except Exception as e:
                    print(f"  ⚠️ 添加截图失败: {e}")

            # 推文文字
            text = tweet.get('text', '').replace('\n', '<br/>')
            if text:
                story.append(self.Paragraph(text, content_style))
                story.append(self.Spacer(1, 0.2 * self.inch))

            # 推文图片
            for img_path in tweet.get('downloaded_images', []):
                if os.path.exists(img_path):
                    try:
                        img = self.Image(img_path, width=5*self.inch, height=3.5*self.inch)
                        story.append(img)
                        story.append(self.Spacer(1, 0.2 * self.inch))
                    except Exception as e:
                        print(f"  ⚠️ 添加图片失败: {e}")

            # 互动数据
            metrics = tweet.get('metrics', {})
            metrics_text = f"💬 {metrics.get('replies', 0)} 回复 · " \
                          f"🔄 {metrics.get('retweets', 0)} 转发 · " \
                          f"❤️ {metrics.get('likes', 0)} 点赞 · " \
                          f"🔖 {metrics.get('bookmarks', 0)} 收藏"
            story.append(self.Paragraph(metrics_text, metrics_style))

            # 分页
            if i < len(data.get('tweets', [])):
                story.append(self.PageBreak())

        # 生成 PDF
        doc.build(story)
        print(f"✅ PDF 生成成功: {output_path}")

        return output_path


class PPTGenerator:
    """PPT 文档生成器"""

    def __init__(self):
        """初始化 PPT 生成器"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor

            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.PP_ALIGN = PP_ALIGN
            self.RGBColor = RGBColor
        except ImportError:
            print("⚠️ python-pptx 未安装，正在安装...")
            import subprocess
            subprocess.run(['pip3', 'install', 'python-pptx'], check=True)
            # 重新导入
            self.__init__()

    def generate(self, data: Dict, output_path: str) -> str:
        """
        生成 PPT 文档

        Args:
            data: Twitter thread 数据
            output_path: 输出文件路径

        Returns:
            生成的文件路径
        """
        print("📊 正在生成 PPT...")

        prs = self.Presentation()
        prs.slide_width = self.Inches(10)
        prs.slide_height = self.Inches(7.5)

        # 获取空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局

        # 封面页
        slide = prs.slides.add_slide(blank_layout)

        # 标题
        author = data.get('author', {})
        author_name = author.get('name', 'Twitter User')
        title_box = slide.shapes.add_textbox(
            self.Inches(1), self.Inches(2),
            self.Inches(8), self.Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{author_name} 的 Twitter Thread"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = self.Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.RGBColor(29, 161, 242)  # Twitter blue
        title_para.alignment = self.PP_ALIGN.CENTER

        # 统计信息
        stats_box = slide.shapes.add_textbox(
            self.Inches(1), self.Inches(3.5),
            self.Inches(8), self.Inches(0.5)
        )
        stats_frame = stats_box.text_frame
        stats_frame.text = f"共 {data.get('total_tweets', 0)} 条推文"
        stats_para = stats_frame.paragraphs[0]
        stats_para.font.size = self.Pt(18)
        stats_para.alignment = self.PP_ALIGN.CENTER

        # 生成时间
        time_box = slide.shapes.add_textbox(
            self.Inches(1), self.Inches(6),
            self.Inches(8), self.Inches(0.5)
        )
        time_frame = time_box.text_frame
        time_frame.text = f"生成时间：{datetime.now().strftime('%Y年%m月%d日')}"
        time_para = time_frame.paragraphs[0]
        time_para.font.size = self.Pt(12)
        time_para.font.color.rgb = self.RGBColor(101, 119, 134)
        time_para.alignment = self.PP_ALIGN.CENTER

        # 为每条推文创建一页
        for i, tweet in enumerate(data.get('tweets', []), 1):
            slide = prs.slides.add_slide(blank_layout)

            # 左侧：推文截图（如果有）
            screenshot = tweet.get('screenshot')
            if screenshot and os.path.exists(screenshot):
                try:
                    slide.shapes.add_picture(
                        screenshot,
                        self.Inches(0.5), self.Inches(1),
                        width=self.Inches(4.5)
                    )
                except Exception as e:
                    print(f"  ⚠️ 添加截图失败: {e}")

            # 右侧：文字内容
            right_left = self.Inches(5.2)
            right_width = self.Inches(4.3)

            # 标题
            title_box = slide.shapes.add_textbox(
                right_left, self.Inches(1),
                right_width, self.Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"推文 #{i}"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = self.Pt(20)
            title_para.font.bold = True

            # 文字内容
            text = tweet.get('text', '')
            if text:
                content_box = slide.shapes.add_textbox(
                    right_left, self.Inches(1.7),
                    right_width, self.Inches(4)
                )
                content_frame = content_box.text_frame
                content_frame.text = text
                content_frame.word_wrap = True
                content_para = content_frame.paragraphs[0]
                content_para.font.size = self.Pt(14)
                content_para.line_spacing = 1.3

            # 互动数据
            metrics = tweet.get('metrics', {})
            metrics_text = f"💬 {metrics.get('replies', 0)} · " \
                          f"🔄 {metrics.get('retweets', 0)} · " \
                          f"❤️ {metrics.get('likes', 0)}"

            metrics_box = slide.shapes.add_textbox(
                right_left, self.Inches(6.5),
                right_width, self.Inches(0.3)
            )
            metrics_frame = metrics_box.text_frame
            metrics_frame.text = metrics_text
            metrics_para = metrics_frame.paragraphs[0]
            metrics_para.font.size = self.Pt(11)
            metrics_para.font.color.rgb = self.RGBColor(101, 119, 134)

        # 保存 PPT
        prs.save(output_path)
        print(f"✅ PPT 生成成功: {output_path}")

        return output_path


def generate_document(
    data: Dict,
    output_path: str,
    format: str = 'pdf'
) -> str:
    """
    生成文档

    Args:
        data: Twitter thread 数据
        output_path: 输出文件路径
        format: 文档格式 ('pdf' 或 'ppt')

    Returns:
        生成的文件路径
    """
    if format.lower() == 'ppt' or format.lower() == 'pptx':
        generator = PPTGenerator()
    else:
        generator = PDFGenerator()

    return generator.generate(data, output_path)
